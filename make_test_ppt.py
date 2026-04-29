import os
from pptx import Presentation
from pptx.util import Inches, Pt


def create_lecture_slide(filename, title, content_list):
    """지정한 제목과 내용으로 PPTX 파일을 생성하는 함수"""
    prs = Presentation()

    # 슬라이드 레이아웃 (제목 + 내용)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    # 제목 설정
    title_shape = slide.shapes.title
    title_shape.text = title

    # 본문 설정
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_list):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = item
        p.level = 0
        p.space_after = Pt(10)

    # 파일 저장
    prs.save(filename)
    print(f"✅ 생성 완료: {filename}")


# --- 전공별 데이터 정의 ---

test_data = {
    "Marketing_Test.pptx": [
        "The AIDA Model in Digital Context",
        "- Attention: Breaking through the noise with targeted social media ads.",
        "- Interest: Providing value-driven content like blog posts.",
        "- Desire: Leveraging social proof and reviews to build brand affinity.",
        "- Action: Optimizing the conversion rate (CRO) on the landing page."
    ],
    "Business_Test.pptx": [
        "Environmental Scanning for Strategy Formulation",
        "- Internal Analysis: Evaluating Core Competencies and Value Chain activities.",
        "- External Analysis: Assessing market trends using the PESTEL framework.",
        "- Strategic Alignment: Ensuring organizational resources match market opportunities."
    ],
    "Economics_Test.pptx": [
        "Determinants of Price Elasticity of Demand",
        "- Availability of Substitutes: If close substitutes are available, demand is more elastic.",
        "- Time Horizon: Demand is more elastic in the long run than in the short run.",
        "- Market Definition: Narrowly defined markets have more elastic demand."
    ],
    "Nursing_Test.pptx": [
        "Standard Precautions and Infection Control",
        "- Hand Hygiene: Most effective way to prevent the spread of microorganisms.",
        "- Personal Protective Equipment (PPE): Proper sequence for donning gloves.",
        "- Aseptic Technique: Maintaining a sterile field during invasive procedures."
    ],
    "Medical_Test.pptx": [
        "Mechanics of the Cardiac Cycle",
        "- Systole: The phase when the heart muscle contracts and pumps blood.",
        "- Diastole: The phase when the heart muscle relaxes and fills with blood.",
        "- Cardiac Output: The volume of blood pumped by the heart per minute (CO = HR x SV)."
    ]
}

# --- 파일 생성 실행 ---
if __name__ == "__main__":
    for filename, content in test_data.items():
        title = content[0]
        body = content[1:]
        create_lecture_slide(filename, title, body)

    print("\n🚀 모든 테스트 파일이 준비되었습니다. EduTrans에서 번역을 시작해보세요!")
import streamlit as st
import fitz  # PyMuPDF
import os
import re
import json
import unicodedata
import tempfile
import subprocess
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from langchain_ollama import OllamaLLM
import platform
from glossary import DEFAULT_GLOSSARY

# --- [1. 초기 설정] ---
st.set_page_config(page_title="TransSlide AI - 로컬 전공 번역기", layout="wide", page_icon="🎓")

USER_GLOSSARY_PATH = "user_glossary.json"


def load_user_glossary():
    if os.path.exists(USER_GLOSSARY_PATH):
        try:
            with open(USER_GLOSSARY_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {}


def save_user_glossary(data):
    with open(USER_GLOSSARY_PATH, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def get_merged_glossary(dept, user_glossary):
    base = dict(DEFAULT_GLOSSARY.get(dept, {}))
    base.update(user_glossary.get(dept, {}))
    return base


# 기본 사전과 사용자 사전을 병합하는 로직 강화
def get_merged_glossary(dept, user_glossary):
    base = dict(DEFAULT_GLOSSARY.get(dept, {}))
    base.update(user_glossary.get(dept, {}))
    return base

# 세션 상태 초기화
if "translation_cache" not in st.session_state:
    st.session_state["translation_cache"] = {}
if "detected_dept" not in st.session_state:
    st.session_state["detected_dept"] = None
if "user_glossary" not in st.session_state:
    st.session_state["user_glossary"] = load_user_glossary()


def get_system_font():
    os_name = platform.system()
    if os_name == "Windows":
        return "C:/Windows/Fonts/malgun.ttf", "Malgun Gothic"
    elif os_name == "Darwin":
        return "/System/Library/Fonts/Supplemental/AppleGothic.ttf", "AppleGothic"
    else:
        return "/usr/share/fonts/truetype/nanum/NanumGothic.ttf", "NanumGothic"


FONT_FILE_PATH, FONT_NAME = get_system_font()


# --- [모델 정보] ---
MODEL_INFO = {
    "qwen2.5:3b": {
        "label": "qwen2.5:3b",
        "ram": "~2GB RAM",
        "quality": "번역 품질 ★★★★★",
        "speed": "속도 ★★★★☆",
        "desc": "Alibaba Qwen2.5 3B. 저사양 환경에서 최고 번역 품질. 강력 추천.",
        "family": "qwen",
        "num_predict": 150,
        "stop": ["\nEnglish:", "\n"],
    },
    "llama3.2:3b": {
        "label": "llama3.2:3b",
        "ram": "~2GB RAM",
        "quality": "번역 품질 ★★★★☆",
        "speed": "속도 ★★★★★",
        "desc": "Meta Llama 3.2 3B. 가장 가볍고 빠름. 번역 품질도 충분히 좋음.",
        "family": "llama",
        "num_predict": 120,
        "stop": ["\nEnglish:", "\n"],
    },
}


# --- [2. 사이드바 UI] ---
# 코드 블록 판별 로직 강화 (어셈블리, C언어 패턴 추가)
def is_code_block(text):
    code_patterns = [
        r'#include', r'\bint\s+\w+\s*\(', r'\bvoid\s+\w+\s*\(',
        r'\bfor\s*\(', r'\bwhile\s*\(', r'\breturn\s+',
        r'printf\s*\(', r'^\s*//', r'[{};]$',
        r'->\w+', r'0x[0-9a-fA-F]+', r'\bmov\w*\s+%',
    ]
    lines = [l for l in text.split('\n') if l.strip()]
    if not lines: return False
    hits = sum(1 for line in lines if any(re.search(p, line) for p in code_patterns))
    return hits / len(lines) > 0.3  # 30% 이상 코드 패턴이면 코드 블록으로 간주


# 과해석 및 사족 제거를 위한 후처리 로직 강화
def post_process(translated, original):
    translated = re.sub(r"<think>.*?</think>", "", translated, flags=re.DOTALL).strip()
    trash_phrases = [
        r"^here is the translation:?", r"^translated text:?", r"^korean:?",
        r"^번역\s*:", r"^결과\s*:", r"^해석\s*:", r"^output\s*:"
    ]
    for phrase in trash_phrases:
        translated = re.sub(phrase, "", translated, flags=re.IGNORECASE).strip()

    # 원문과 똑같거나 비어있으면 무시
    if not translated or translated.lower() == original.lower():
        return ""
    return translated


with st.sidebar:
    st.header("⚙️ 엔진 설정")
    selected_model = st.selectbox(
        "사용할 AI 모델 선택",
        list(MODEL_INFO.keys()),
        index=0,
        format_func=lambda m: MODEL_INFO[m]["label"],
    )

    info = MODEL_INFO[selected_model]
    st.caption(f"{info['ram']} | {info['quality']} | {info['speed']}\n\n{info['desc']}")

    llm = OllamaLLM(
        model=selected_model,
        temperature=0,
        num_predict=MODEL_INFO[selected_model]["num_predict"],
        timeout=60,  # 60초 타임아웃 (응답 없음 방지)
    )

    st.divider()
    st.header("🏫 전공 및 과목 설정")

    auto_detect = st.checkbox(
        "🔍 전공 자동 감지 모드", value=True,
        help="파일 첫 페이지를 분석해 전공을 자동으로 추론합니다. 비활성화하면 아래 수동 선택이 적용됩니다."
    )

    category = st.radio(
        "수동 전공 선택 (자동 감지 비활성화 시 사용)",
        ["문과 (Humanities)", "이과 (Science & Engineering)", "예체능 (Arts & Sports)"],
        disabled=auto_detect,
    )

    dept_map = {
        "문과 (Humanities)": [
            "Business Administration", "Economics", "Marketing", "Accounting",
            "International Relations", "Psychology", "Sociology",
            "Media & Communication", "English Language & Literature", "Political Science"
        ],
        "이과 (Science & Engineering)": [
            "Computer Science", "Data Science", "Artificial Intelligence Engineering",
            "Electrical Engineering", "Mechanical Engineering", "Chemical Engineering",
            "Biotechnology", "Nursing", "Architecture", "Civil Engineering",
            "Mathematics", "Physics", "Medical"
        ],
        "예체능 (Arts & Sports)": [
            "Fine Arts", "Graphic Design", "Industrial Design", "Music Composition",
            "Vocal Music", "Physical Education", "Film & Digital Media", "Fashion Design"
        ]
    }
    manual_dept = st.selectbox(
        "세부 전공", dept_map[category],
        disabled=auto_detect,
    )

    if auto_detect:
        detected = st.session_state.get("detected_dept")
        if detected and detected != "General":
            st.info(f"📍 감지된 전공: **{detected}**")
        else:
            st.info("📍 파일 업로드 시 전공을 자동 감지합니다.")
    else:
        st.info(f"📍 수동 선택 전공: **{manual_dept}**")

    st.divider()
    st.header("📚 용어 사전 편집")

    active_dept = (st.session_state.get("detected_dept") or manual_dept) if auto_detect else manual_dept
    st.caption(f"현재 전공: **{active_dept}**")

    with st.expander("기본 용어 사전 보기"):
        default_terms = DEFAULT_GLOSSARY.get(active_dept, {})
        if default_terms:
            for eng, kor in default_terms.items():
                st.markdown(f"- **{eng}** → {kor}")
        else:
            st.caption("이 전공의 기본 용어가 아직 없습니다.")

    with st.expander("✏️ 내 용어 추가 / 수정"):
        with st.form("add_term_form", clear_on_submit=True):
            col1, col2 = st.columns(2)
            with col1:
                new_eng = st.text_input("영어 용어", placeholder="e.g. Gradient Descent")
            with col2:
                new_kor = st.text_input("한국어 번역", placeholder="e.g. 경사 하강법")
            submitted = st.form_submit_button("➕ 추가")
            if submitted and new_eng.strip() and new_kor.strip():
                if active_dept not in st.session_state["user_glossary"]:
                    st.session_state["user_glossary"][active_dept] = {}
                st.session_state["user_glossary"][active_dept][new_eng.strip()] = new_kor.strip()
                save_user_glossary(st.session_state["user_glossary"])
                st.success(f"추가됨: {new_eng.strip()} → {new_kor.strip()}")

    user_dept_terms = st.session_state["user_glossary"].get(active_dept, {})
    if user_dept_terms:
        with st.expander(f"내 커스텀 용어 목록 ({len(user_dept_terms)}개)"):
            to_delete = []
            for eng, kor in user_dept_terms.items():
                col_t, col_d = st.columns([4, 1])
                with col_t:
                    st.markdown(f"**{eng}** → {kor}")
                with col_d:
                    if st.button("🗑️", key=f"del_{eng}"):
                        to_delete.append(eng)
            for key in to_delete:
                del st.session_state["user_glossary"][active_dept][key]
                save_user_glossary(st.session_state["user_glossary"])
                st.rerun()

    with st.expander("📤 용어 사전 내보내기 / 가져오기"):
        export_data = json.dumps(st.session_state["user_glossary"], ensure_ascii=False, indent=2)
        st.download_button(
            "💾 내 용어 사전 다운로드 (JSON)",
            data=export_data,
            file_name="my_glossary.json",
            mime="application/json"
        )
        uploaded_glossary = st.file_uploader("JSON 용어 사전 불러오기", type="json", key="glossary_upload")
        if uploaded_glossary:
            try:
                imported = json.loads(uploaded_glossary.read().decode("utf-8"))
                for dept, terms in imported.items():
                    if dept not in st.session_state["user_glossary"]:
                        st.session_state["user_glossary"][dept] = {}
                    st.session_state["user_glossary"][dept].update(terms)
                save_user_glossary(st.session_state["user_glossary"])
                st.success("용어 사전을 성공적으로 가져왔습니다.")
            except Exception as e:
                st.error(f"파일 읽기 오류: {e}")

    st.divider()
    if st.button("🧹 번역 캐시 초기화"):
        st.session_state["translation_cache"] = {}
        st.session_state["detected_dept"] = None
        st.success("캐시가 초기화되었습니다.")


# --- [3. 핵심 유틸리티] ---

def get_model_family():
    return MODEL_INFO.get(selected_model, {}).get("family", "llama")


def normalize_for_cache(text):
    return re.sub(r"\s+", " ", text.lower().strip())


def clean_text_logic(text):
    text = unicodedata.normalize("NFKC", text)
    text = text.replace("-\n", "").replace("\n", " ")
    return re.sub(r"\s+", " ", text).strip()


def is_english_content(text):
    ascii_alpha = sum(1 for c in text if c.isalpha() and ord(c) < 128)
    korean = sum(1 for c in text if '가' <= c <= '힣')
    total_alpha = sum(1 for c in text if c.isalpha())
    if total_alpha == 0 or ascii_alpha == 0:
        return False
    return korean / total_alpha < 0.3


def is_code_block(text):
    code_patterns = [
        r'#include', r'\bint\s+\w+\s*\(', r'\bvoid\s+\w+\s*\(',
        r'\bfor\s*\(', r'\bwhile\s*\(', r'\breturn\s+',
        r'printf\s*\(', r'^\s*//', r'[{};]$',
        r'->\w+', r'0x[0-9a-fA-F]+',
        r'\bmov\w*\s+%', r'\bpush\b', r'\bpop\b',
        r'\.text\b', r'\.data\b',
    ]
    lines = [l for l in text.split('\n') if l.strip()]
    if len(lines) < 2:
        return False
    hits = sum(1 for line in lines if any(re.search(p, line) for p in code_patterns))
    return hits / len(lines) > 0.35


def detect_major_from_text(text, fallback_dept):
    if not text.strip():
        return fallback_dept
    prompt = (
        f"Read this university lecture slide text and identify the academic major/subject.\n"
        f"Reply with ONLY the major name in English (e.g. 'Data Science', 'Nursing', 'Marketing').\n"
        f"Do NOT explain. Output the major name only.\n\n"
        f"Text: {text[:600]}\nMajor:"
    )
    try:
        result = str(llm.invoke(prompt, stop=["\n", "\n\n"])).strip()
        result = re.sub(r"<think>.*?</think>", "", result, flags=re.DOTALL).strip()
        result = result.split("\n")[0].strip()
        if not result or result.lower() in ("general", "unknown", "none", "n/a"):
            return fallback_dept
        return result
    except Exception:
        return fallback_dept


def glossary_exact_match(text, glossary):
    stripped = text.strip().rstrip(":")
    if stripped in glossary:
        return glossary[stripped]
    lower = stripped.lower()
    for k, v in glossary.items():
        if k.lower() == lower:
            return v
    return None


def build_glossary_hint(glossary, text):
    if not glossary:
        return ""
    text_lower = text.lower()
    matched = {k: v for k, v in glossary.items() if k.lower() in text_lower}
    if not matched:
        return ""
    items = list(matched.items())[:6]
    pairs = ", ".join(f"{k}={v}" for k, v in items)
    return f"[Terms: {pairs}]"


def post_process(translated, original):
    translated = re.sub(r"<think>.*?</think>", "", translated, flags=re.DOTALL).strip()
    trash_phrases = [
        r"^here is the translation:?", r"^translated text:?", r"^korean translation:?",
        r"^번역\s*:", r"^결과\s*:", r"^해석\s*:", r"^output\s*:", r"^korean\s*:",
        r"^강의\s*:", r"^강의\s+\d+\s*:", r"^translation\s*:", r"^/no_think",
    ]
    for phrase in trash_phrases:
        translated = re.sub(phrase, "", translated, flags=re.IGNORECASE).strip()
    translated = re.sub(r"[^가-힣 -~ -ɏ　-〿().,!?:/\-\d]", "", translated).strip()
    translated = re.sub(r'^["\' ]+|["\' ]+$', "", translated).strip()
    lines = [l.strip() for l in translated.split("\n") if l.strip()]
    if lines:
        translated = lines[0] if len(lines[0]) > 1 else " ".join(lines[:2])
    if not translated or translated.lower() == original.lower():
        return ""
    return translated


def translate_single(text, department):
    cleaned = clean_text_logic(text)
    if len(cleaned) < 1:
        return text

    norm_key = normalize_for_cache(cleaned)
    if norm_key in st.session_state["translation_cache"]:
        return st.session_state["translation_cache"][norm_key]

    glossary = get_merged_glossary(department, st.session_state["user_glossary"])

    if len(cleaned.split()) <= 6:
        exact = glossary_exact_match(cleaned, glossary)
        if exact:
            st.session_state["translation_cache"][norm_key] = exact
            return exact

    family = get_model_family()
    glossary_hint = build_glossary_hint(glossary, cleaned)
    stop_param = MODEL_INFO[selected_model]["stop"]

    if family == "qwen":
        prompt = (
            f"You are a Korean translator for {department} academic slides.\n"
            f"Translate the English text into Korean. ONE LINE only. No lists. No explanations.\n"
            f"{glossary_hint}\n\n"
            f"English: {cleaned}\n"
            f"Korean:"
        )
    else:
        prompt = (
            f"You are a Korean translator for {department} academic slides.\n"
            f"Translate the English text into Korean. ONE SHORT LINE only. No lists. No explanations. No extra sentences.\n"
            f"{glossary_hint}\n\n"
            f"English: {cleaned}\n"
            f"Korean:"
        )

    try:
        translated = str(llm.invoke(prompt, stop=stop_param)).strip()
        translated = post_process(translated, cleaned)
        if not translated:
            return ""
        st.session_state["translation_cache"][norm_key] = translated
        return translated
    except Exception as e:
        st.error(f"Ollama 연결 확인 필요: {e}")
        return text


# --- [4. 문서 처리 함수] ---

def resolve_dept(auto_detect_flag, fallback_dept, first_text):
    if not auto_detect_flag:
        return fallback_dept
    detected = detect_major_from_text(first_text, fallback_dept)
    st.session_state["detected_dept"] = detected
    return detected


def process_pptx(input_path, output_path, fallback_dept, auto_detect_flag):
    prs = Presentation(input_path)
    total_slides = len(prs.slides)
    progress_bar = st.progress(0)
    status_text = st.empty()

    first_text = ""
    if total_slides > 0:
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text"):
                first_text += shape.text + " "

    dept = resolve_dept(auto_detect_flag, fallback_dept, first_text)
    status_text.text(f"📌 적용 전공: {dept}")

    for i, slide in enumerate(prs.slides):
        status_text.text(f"슬라이드 {i + 1}/{total_slides} 번역 중... (전공: {dept})")
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    original_text = "".join(run.text for run in paragraph.runs).strip()
                    if len(original_text) < 2 or not is_english_content(original_text):
                        continue
                    translated = translate_single(original_text, dept)
                    if translated and translated != original_text:
                        run = paragraph.add_run()
                        run.text = f"\n{translated}"
                        run.font.name = FONT_NAME
                        run.font.size = Pt(10)
                        run.font.color.rgb = RGBColor(0, 102, 204)
        progress_bar.progress((i + 1) / total_slides)

    prs.save(output_path)
    status_text.text(f"✅ PPT 번역 완료! (전공: {dept})")


def process_pdf(input_path, output_path, fallback_dept, auto_detect_flag):
    doc = fitz.open(input_path)
    total_pages = len(doc)
    progress_bar = st.progress(0)
    status_text = st.empty()

    first_text = ""
    if total_pages > 0:
        first_text = " ".join(
            [b[4] for b in doc[0].get_text("blocks") if b[4].strip()]
        )

    dept = resolve_dept(auto_detect_flag, fallback_dept, first_text)
    status_text.text(f"📌 적용 전공: {dept}")

    for i, page in enumerate(doc):
        status_text.text(f"📄 페이지 {i + 1}/{total_pages} 번역 중... (전공: {dept})")
        blocks = page.get_text("blocks")
        page_height = page.rect.height

        for b in blocks:
            raw_text = b[4].strip()
            if len(raw_text) < 3 or not is_english_content(raw_text):
                continue

            if is_code_block(raw_text):
                continue
            cleaned_text = raw_text.replace("\n", " ").strip()
            translated = translate_single(cleaned_text, dept)

            if translated and translated != cleaned_text:
                insert_x = b[0]
                insert_y = b[3] + 7
                if insert_y > page_height - 10:
                    continue
                try:
                    page.insert_text(
                        (insert_x, insert_y),
                        translated,
                        fontname="ko",
                        fontfile=FONT_FILE_PATH,
                        fontsize=8,
                        color=(0, 0.4, 0.8)
                    )
                except Exception:
                    continue

        progress_bar.progress((i + 1) / total_pages)

    doc.save(output_path)
    doc.close()
    status_text.text(f"✅ PDF 번역 완료! (전공: {dept})")


# --- [5. 메인 UI] ---
st.markdown("""<style>
    .stSidebar h2 { font-size: 1.8rem !important; color: #4A90E2; }
    .stButton button { width: 100%; background-color: #4A90E2 !important; color: white !important; border-radius: 10px; }
</style>""", unsafe_allow_html=True)

st.title("🎓 전공 맞춤형 강의자료 번역기")

with st.expander("📖 이용 가이드 및 주의사항", expanded=False):
    st.markdown("""
    1. **로컬 엔진:** Ollama가 실행 중이어야 합니다 (`ollama serve`).
    2. **모델 선택:** **qwen2.5:3b** (1순위, 번역 품질 최고) 또는 **llama3.2:3b** (2순위, 속도 최고)를 권장합니다.
    3. **모델 설치:** 앱 시작 시 두 모델이 자동으로 설치됩니다.
    4. **자동 감지 (기본값 ON):** 파일 첫 페이지를 분석해 전공을 자동 추론합니다.
    5. **수동 모드:** 자동 감지를 끄면 사이드바에서 직접 전공을 선택합니다.
    6. **용어 사전:** 사이드바에서 전공별 커스텀 용어를 추가·삭제하고 JSON으로 저장할 수 있습니다.
    7. **보안:** 모든 데이터는 본인 PC 내에서만 처리됩니다.
    """)

tab1, tab2 = st.tabs(["📊 PPT 번역", "📄 PDF 번역"])


def get_save_path(filename):
    downloads = os.path.join(os.path.expanduser("~"), "Downloads")
    os.makedirs(downloads, exist_ok=True)
    return os.path.join(downloads, f"translated_{filename}")


def open_folder(path):
    if platform.system() == "Windows":
        os.startfile(path)
    elif platform.system() == "Darwin":
        subprocess.Popen(["open", path])
    else:
        subprocess.Popen(["xdg-open", path])


def run_translation(uploaded_file, mode):
    result_key = f"result_{mode}"
    file_key = f"file_{mode}"

    # 새 파일 업로드 시 이전 결과 초기화
    if uploaded_file:
        if st.session_state.get(file_key) != uploaded_file.name:
            st.session_state[result_key] = None
            st.session_state[file_key] = uploaded_file.name

    if uploaded_file:
        if st.button(f"🚀 {mode} 번역 시작"):
            suffix = os.path.splitext(uploaded_file.name)[1]
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp_in:
                tmp_in.write(uploaded_file.getvalue())
                in_path = tmp_in.name

            out_filename = f"translated_{uploaded_file.name}"
            out_path = get_save_path(uploaded_file.name)

            try:
                with st.spinner(f"{selected_model} 모델 번역 중..."):
                    if mode == "PPTX":
                        process_pptx(in_path, out_path, manual_dept, auto_detect)
                    else:
                        process_pdf(in_path, out_path, manual_dept, auto_detect)

                st.session_state[result_key] = {
                    "out_path": out_path,
                    "out_filename": out_filename,
                }
            except Exception as e:
                st.error(f"오류 발생: {e}")
            finally:
                if os.path.exists(in_path):
                    os.remove(in_path)

        result = st.session_state.get(result_key)
        if result and os.path.exists(result["out_path"]):
            out_path = result["out_path"]
            out_filename = result["out_filename"]
            st.success("✅ 번역 완료!")
            st.info(f"📁 저장 위치: `{out_path}`")

            col1, col2 = st.columns(2)
            with col1:
                if st.button("📂 저장 폴더 열기", key=f"open_{mode}"):
                    open_folder(os.path.dirname(out_path))
            with col2:
                with open(out_path, "rb") as f:
                    st.download_button("💾 직접 다운로드", f,
                                       file_name=out_filename,
                                       key=f"dl_{mode}")


with tab1:
    run_translation(st.file_uploader("PPTX 업로드", type="pptx", key="p1"), "PPTX")
with tab2:
    run_translation(st.file_uploader("PDF 업로드", type="pdf", key="p2"), "PDF")

st.divider()
st.markdown("""
    <div style="text-align: center; color: gray; font-size: 0.75rem; line-height: 1.8;">
        <p><b>Built with Qwen</b> · <b>Built with Llama</b> · <b>Built with Meta Llama 3</b></p>
        <p>
            <b>Qwen2.5</b>: Qwen Research License (비상업적/연구·평가 목적) © Alibaba Cloud. All Rights Reserved.<br>
            <b>Llama 3.2</b>: Llama 3.2 Community License © Meta Platforms, Inc. All Rights Reserved.<br>
        </p>
        <p style="font-size: 0.7rem; opacity: 0.75;">
            ※ Qwen2.5 모델은 비상업적 용도(연구·교육·평가)로만 사용 가능합니다. 상업적 사용 시 Alibaba Cloud에 별도 라이선스를 요청하세요.<br>
            ※ Llama 3.2 모델은 월간 활성 사용자 7억 명 미만의 서비스에 한해 상업적 사용이 허용됩니다.<br>
            모든 데이터는 로컬 PC에서만 처리됩니다. 외부 서버로 전송되지 않습니다.
        </p>
    </div>
    """, unsafe_allow_html=True)